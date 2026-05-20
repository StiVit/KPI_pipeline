from pptx import Presentation

def create_ppt(kpis, image_paths):
    """
    This is an example of a pptx generator, it needs to be changed completely to generate the pptx that we need

    ARGS:
        - kpis: The kpis measured before
        - image_paths: All the plots made from the folder

    TODO: Create customisable settings for the pptx
    FIXME: Change the whole function to generate the required pptx based on the settings
    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout=slide_layout)
    slide.shapes.title.text = "KPI presentatione"

    # KPI slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout=slide_layout)

    content = "\n".join([f"{k}: {v}" for k, v in kpis.items()])
    slide.placeholders[1].text = content

    # Chart slide
    for img in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img, 100, 100)

    prs.save()